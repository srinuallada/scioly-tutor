"""
Document processor — extracts and chunks text from Word, PowerPoint, PDF, Excel, and text files.

Usage:
    python -m retrieval.processor /path/to/materials/folder
    python -m retrieval.processor file1.docx file2.pptx
"""

import json
import hashlib
import re
import os
import sys
from pathlib import Path
from dataclasses import dataclass, asdict, field
from typing import Optional


@dataclass
class Chunk:
    """A searchable piece of study material."""
    source_file: str
    source_type: str
    section_title: str
    content: str
    page_or_slide: Optional[int] = None
    word_count: int = 0
    id: str = ""

    def __post_init__(self):
        self.word_count = len(self.content.split())
        if not self.id:
            raw = f"{self.source_file}:{self.section_title}:{self.content[:100]}"
            self.id = hashlib.md5(raw.encode()).hexdigest()[:12]


# ──────────────────────────────────────────────
# EXTRACTORS
# ──────────────────────────────────────────────

def extract_docx(filepath: str) -> list[Chunk]:
    """Extract text from Word documents, chunked by heading."""
    from docx import Document

    doc = Document(filepath)
    fname = Path(filepath).name
    chunks: list[Chunk] = []
    current_heading = "Introduction"
    current_text: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        if para.style.name.startswith("Heading"):
            if current_text:
                chunks.append(Chunk(
                    source_file=fname,
                    source_type="docx",
                    section_title=current_heading,
                    content="\n".join(current_text),
                ))
            current_heading = text
            current_text = []
        else:
            current_text.append(text)

    if current_text:
        chunks.append(Chunk(
            source_file=fname,
            source_type="docx",
            section_title=current_heading,
            content="\n".join(current_text),
        ))

    # Extract tables
    for i, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        if rows:
            chunks.append(Chunk(
                source_file=fname,
                source_type="docx",
                section_title=f"Table {i + 1}",
                content="\n".join(rows),
            ))

    return chunks


def extract_pptx(filepath: str) -> list[Chunk]:
    """Extract text from PowerPoint, one chunk per slide + speaker notes."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(filepath)
    fname = Path(filepath).name
    chunks: list[Chunk] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        title = ""
        body_parts: list[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                if (hasattr(shape, "placeholder_format")
                        and shape.placeholder_format is not None
                        and shape.placeholder_format.idx == 0):
                    title = shape.text_frame.text.strip()
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        body_parts.append(text)

            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    body_parts.append(" | ".join(cells))

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                body_parts.append("[Image on slide]")
            if shape.has_chart:
                body_parts.append("[Chart on slide]")

        if not title:
            title = f"Slide {slide_num}"

        content = "\n".join(body_parts)
        if content.strip():
            chunks.append(Chunk(
                source_file=fname,
                source_type="pptx",
                section_title=title,
                content=content,
                page_or_slide=slide_num,
            ))

        # Speaker notes often contain extra study info
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunks.append(Chunk(
                    source_file=fname,
                    source_type="pptx",
                    section_title=f"Notes — {title}",
                    content=notes,
                    page_or_slide=slide_num,
                ))

    return chunks


def extract_pdf(filepath: str) -> list[Chunk]:
    """Extract text from PDF, chunked by page."""
    import pdfplumber

    fname = Path(filepath).name
    chunks: list[Chunk] = []

    with pdfplumber.open(filepath) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                chunks.append(Chunk(
                    source_file=fname,
                    source_type="pdf",
                    section_title=f"Page {page_num}",
                    content=text.strip(),
                    page_or_slide=page_num,
                ))

            tables = page.extract_tables()
            for t_idx, table in enumerate(tables):
                if table:
                    rows = []
                    for row in table:
                        cells = [str(c).strip() if c else "" for c in row]
                        rows.append(" | ".join(cells))
                    chunks.append(Chunk(
                        source_file=fname,
                        source_type="pdf",
                        section_title=f"Page {page_num} — Table {t_idx + 1}",
                        content="\n".join(rows),
                        page_or_slide=page_num,
                    ))

    return chunks


def extract_xlsx(filepath: str) -> list[Chunk]:
    """Extract text from Excel, one chunk per sheet."""
    from openpyxl import load_workbook

    wb = load_workbook(filepath, data_only=True)
    fname = Path(filepath).name
    chunks: list[Chunk] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            line = " | ".join(cells)
            if line.replace("|", "").strip():
                rows.append(line)

        if rows:
            chunks.append(Chunk(
                source_file=fname,
                source_type="xlsx",
                section_title=sheet_name,
                content="\n".join(rows),
            ))

    return chunks


def extract_text(filepath: str) -> list[Chunk]:
    """Extract from plain text / markdown files."""
    fname = Path(filepath).name
    ext = Path(filepath).suffix.lower()

    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        content = f.read()

    if not content.strip():
        return []

    if ext in (".md", ".markdown"):
        sections = re.split(r'\n(#{1,3}\s+.+)', content)
        chunks: list[Chunk] = []
        current_title = "Introduction"

        for part in sections:
            part = part.strip()
            if not part:
                continue
            if re.match(r'^#{1,3}\s+', part):
                current_title = re.sub(r'^#{1,3}\s+', '', part)
            else:
                chunks.append(Chunk(
                    source_file=fname,
                    source_type="md",
                    section_title=current_title,
                    content=part,
                ))
        return chunks if chunks else [Chunk(
            source_file=fname, source_type="md",
            section_title="Full Document", content=content,
        )]

    # Plain text: chunk by double newlines
    paragraphs = re.split(r'\n\s*\n', content)
    chunks = []
    for i, para in enumerate(paragraphs):
        para = para.strip()
        if para and len(para.split()) > 5:
            chunks.append(Chunk(
                source_file=fname,
                source_type="txt",
                section_title=f"Section {i + 1}",
                content=para,
            ))

    return chunks if chunks else [Chunk(
        source_file=fname, source_type="txt",
        section_title="Full Document", content=content,
    )]


# ──────────────────────────────────────────────
# MAIN PROCESSOR
# ──────────────────────────────────────────────

EXTRACTORS = {
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".pdf": extract_pdf,
    ".xlsx": extract_xlsx,
    ".xls": extract_xlsx,
    ".txt": extract_text,
    ".md": extract_text,
    ".markdown": extract_text,
    ".csv": extract_text,
}


def process_file(filepath: str) -> list[Chunk]:
    """Process a single file and return chunks."""
    ext = Path(filepath).suffix.lower()
    extractor = EXTRACTORS.get(ext)

    if not extractor:
        print(f"  ⚠️  Unsupported: {ext} ({filepath})")
        return []

    try:
        chunks = extractor(filepath)
        print(f"  ✅ {Path(filepath).name}: {len(chunks)} chunks")
        return chunks
    except Exception as e:
        print(f"  ❌ Error {filepath}: {e}")
        return []


def process_directory(directory: str) -> list[Chunk]:
    """Process all supported files in a directory (recursive)."""
    all_chunks: list[Chunk] = []
    dir_path = Path(directory)

    if not dir_path.exists():
        print(f"Directory not found: {directory}")
        return []

    supported_files = []
    for ext in EXTRACTORS:
        supported_files.extend(dir_path.rglob(f"*{ext}"))

    print(f"\n📂 Found {len(supported_files)} files in {directory}\n")

    for filepath in sorted(supported_files):
        chunks = process_file(str(filepath))
        all_chunks.extend(chunks)

    print(f"\n📊 Total: {len(all_chunks)} chunks, {sum(c.word_count for c in all_chunks):,} words")
    return all_chunks


def save_chunks(chunks: list[Chunk], output_path: str = "chunks.json") -> None:
    """Save processed chunks to JSON."""
    data = [asdict(c) for c in chunks]
    with open(output_path, "w") as f:
        json.dump(data, f, indent=2)
    print(f"💾 Saved {len(chunks)} chunks to {output_path}")


def load_chunks(input_path: str = "chunks.json") -> list[dict]:
    """Load chunks from JSON as dicts."""
    with open(input_path, "r") as f:
        return json.load(f)


# ──────────────────────────────────────────────
# CLI
# ──────────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage:")
        print("  python -m retrieval.processor /path/to/materials/")
        print("  python -m retrieval.processor file1.docx file2.pptx")
        print(f"\nSupported: {', '.join(sorted(EXTRACTORS.keys()))}")
        sys.exit(1)

    target = sys.argv[1]

    # Determine output path
    data_dir = os.path.join(os.path.dirname(__file__), "..", "..", "data")
    os.makedirs(data_dir, exist_ok=True)
    chunks_path = os.path.join(data_dir, "chunks.json")

    if os.path.isdir(target):
        chunks = process_directory(target)
    else:
        chunks = []
        for fp in sys.argv[1:]:
            chunks.extend(process_file(fp))
        print(f"\n📊 Total: {len(chunks)} chunks")

    if chunks:
        save_chunks(chunks, chunks_path)
