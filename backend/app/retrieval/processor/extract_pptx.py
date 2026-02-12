"""Extract text from PowerPoint, one chunk per slide + speaker notes."""

import logging
from pathlib import Path
from urllib.parse import quote

from app.domain.documents import Chunk
from app.retrieval.processor.chunking_utils import split_text_to_chunks
from app.settings import IMAGES_DIR

log = logging.getLogger(__name__)

# Map common MIME content types to file extensions
_CONTENT_TYPE_TO_EXT = {
    "image/png": "png",
    "image/jpeg": "jpg",
    "image/gif": "gif",
    "image/bmp": "bmp",
    "image/tiff": "tiff",
    "image/svg+xml": "svg",
    "image/x-emf": "emf",
    "image/x-wmf": "wmf",
}


def _save_image(blob: bytes, content_type: str, slide_num: int, img_index: int, stem: str) -> str:
    """Save image blob to IMAGES_DIR and return the filename."""
    ext = _CONTENT_TYPE_TO_EXT.get(content_type, "png")
    filename = f"{stem}_s{slide_num}_{img_index}.{ext}"
    out_path = IMAGES_DIR / filename
    out_path.write_bytes(blob)
    return filename


def extract_pptx(filepath: str) -> list[Chunk]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(filepath)
    fname = Path(filepath).name
    stem = Path(filepath).stem
    chunks: list[Chunk] = []

    IMAGES_DIR.mkdir(parents=True, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides, 1):
        title = ""
        body_parts: list[str] = []
        img_index = 0

        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    title = shape.text_frame.text.strip()
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        body_parts.append(text)

            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    body_parts.append(" | ".join(cells))

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    ct = shape.image.content_type
                    filename = _save_image(blob, ct, slide_num, img_index, stem)
                    body_parts.append(f"![Slide diagram](/api/images/{quote(filename)})")
                    img_index += 1
                except Exception:
                    log.warning("Could not extract image from slide %d", slide_num, exc_info=True)
                    body_parts.append("[Image on slide]")
            if shape.has_chart:
                body_parts.append("[Chart on slide]")

        if not title:
            title = f"Slide {slide_num}"

        content = "\n".join(body_parts)
        if content.strip():
            chunk_texts = split_text_to_chunks(content)
            if not chunk_texts and content:
                chunk_texts = [content]
            for idx, chunk_text in enumerate(chunk_texts, 1):
                chunks.append(Chunk(
                    source_file=fname, source_type="pptx",
                    section_title=f"{title} — Part {idx}",
                    content=chunk_text,
                    page_or_slide=slide_num,
                ))

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                for idx, chunk_text in enumerate(split_text_to_chunks(notes), 1):
                    chunks.append(Chunk(
                        source_file=fname, source_type="pptx",
                        section_title=f"Notes — {title} — Part {idx}",
                        content=chunk_text, page_or_slide=slide_num,
                    ))

    return chunks
